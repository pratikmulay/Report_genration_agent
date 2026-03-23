"""
python-pptx based PPTX exporter.
Generates 16:9 PowerPoint presentations with one slide per section.
Title slide → Executive Summary → Section slides → Key Takeaways.
"""

import base64
import io
import logging
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from app.models import NarrativeSections, BrandingConfig

logger = logging.getLogger(__name__)

# 16:9 dimensions
SLIDE_WIDTH = Inches(13.33)
SLIDE_HEIGHT = Inches(7.5)


def _hex_to_pptx_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to pptx RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _add_title_slide(
    prs: Presentation,
    branding: BrandingConfig,
    user_query: str,
):
    """Add the title slide with branding."""
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = "Analysis Report"
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(44)
            run.font.color.rgb = _hex_to_pptx_rgb(branding.primary_color)
            run.font.bold = True

    # Subtitle with company name and query
    subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if subtitle:
        lines = []
        if branding.company_name:
            lines.append(branding.company_name)
        if user_query:
            lines.append(f"Query: {user_query}")
        subtitle.text = "\n".join(lines)
        for para in subtitle.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(18)

    # Logo
    if branding.logo_base64:
        try:
            logo_bytes = base64.b64decode(branding.logo_base64)
            logo_stream = io.BytesIO(logo_bytes)
            slide.shapes.add_picture(
                logo_stream,
                left=Inches(0.5),
                top=Inches(0.3),
                width=Inches(1.5),
            )
        except Exception as e:
            logger.warning(f"Failed to add logo to title slide: {e}")


def _add_content_slide(
    prs: Presentation,
    title: str,
    content: str,
    accent_color: RGBColor,
    chart: Optional[dict] = None,
):
    """Add a content slide with title, text, and optional chart image."""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(32)
            run.font.color.rgb = accent_color
            run.font.bold = True

    # Content body
    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body:
        # Truncate content for slide readability
        paragraphs = content.split("\n")
        body.text = ""
        tf = body.text_frame
        tf.word_wrap = True

        for i, para_text in enumerate(paragraphs[:6]):  # Max 6 paragraphs
            if i == 0:
                tf.text = para_text.strip()
            else:
                p = tf.add_paragraph()
                p.text = para_text.strip()
            for para in tf.paragraphs:
                para.font.size = Pt(14)

    # Embed chart image if provided
    if chart:
        image_b64 = chart.get("image_base64", chart.get("png_base64", ""))
        if image_b64:
            try:
                img_bytes = base64.b64decode(image_b64)
                img_stream = io.BytesIO(img_bytes)
                slide.shapes.add_picture(
                    img_stream,
                    left=Inches(7.0),
                    top=Inches(1.5),
                    width=Inches(5.5),
                    height=Inches(4.5),
                )
            except Exception as e:
                logger.warning(f"Failed to embed chart on slide: {e}")


def _add_takeaways_slide(
    prs: Presentation,
    narratives: NarrativeSections,
    accent_color: RGBColor,
):
    """Add a key takeaways slide summarizing the report."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = "Key Takeaways"
    for para in title_shape.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(36)
            run.font.color.rgb = accent_color
            run.font.bold = True

    body = slide.placeholders[1] if len(slide.placeholders) > 1 else None
    if body:
        tf = body.text_frame
        tf.word_wrap = True
        takeaways = []

        if narratives.executive_summary:
            # Extract first sentence as key takeaway
            first_sentence = narratives.executive_summary.split(".")[0] + "."
            takeaways.append(f"• {first_sentence.strip()}")

        sections = {
            "Data": narratives.data_overview,
            "SQL": narratives.sql_findings,
            "ML": narratives.ml_insights,
            "NLP": narratives.nlp_section,
        }
        for name, text in sections.items():
            if text:
                first_sentence = text.split(".")[0] + "."
                takeaways.append(f"• {name}: {first_sentence.strip()}")

        tf.text = "\n\n".join(takeaways) if takeaways else "No findings available."
        for para in tf.paragraphs:
            para.font.size = Pt(16)


def export_pptx(
    narratives: NarrativeSections,
    charts: Optional[list[dict]] = None,
    branding: Optional[BrandingConfig] = None,
    report_style: str = "detailed",
    user_query: str = "",
    include_charts: bool = True,
) -> bytes:
    """
    Generate a PPTX presentation.

    Returns:
        PPTX file as bytes
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    brand = branding or BrandingConfig()
    accent = _hex_to_pptx_rgb(brand.primary_color)

    chart_list = charts or []
    chart_idx = 0

    # 1. Title slide
    _add_title_slide(prs, brand, user_query)

    # 2. Executive Summary slide
    if narratives.executive_summary:
        _add_content_slide(prs, "Executive Summary", narratives.executive_summary, accent)

    # 3. Section slides
    sections = [
        ("Data Overview", narratives.data_overview),
        ("SQL Analysis Findings", narratives.sql_findings),
        ("Machine Learning Insights", narratives.ml_insights),
        ("NLP & Text Analysis", narratives.nlp_section),
    ]

    for title, content in sections:
        if content:
            chart = None
            if include_charts and chart_idx < len(chart_list):
                chart = chart_list[chart_idx]
                chart_idx += 1
            _add_content_slide(prs, title, content, accent, chart=chart)

    # 4. Remaining charts on their own slides
    if include_charts:
        while chart_idx < len(chart_list):
            chart = chart_list[chart_idx]
            _add_content_slide(
                prs,
                chart.get("title", "Visualization"),
                chart.get("description", ""),
                accent,
                chart=chart,
            )
            chart_idx += 1

    # 5. Key Takeaways slide
    _add_takeaways_slide(prs, narratives, accent)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    pptx_bytes = buffer.read()
    logger.info(f"Generated PPTX presentation ({len(pptx_bytes)} bytes)")
    return pptx_bytes
